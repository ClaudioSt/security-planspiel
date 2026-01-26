#!/usr/bin/env python3
"""
PowerPoint-Generator mit Template-Unterstützung
Security-Planspiel MechTech

Verwendet bestehende PowerPoint-Templates und füllt sie mit Daten aus simulation_config.json

Templates:
- pptx_output/Massnahmenkarten_Security-Game.pptx (30 Folien: 10 Maßnahmen × 3 Levels)
- pptx_output/Events_Security-Game.pptx (1 Folie als Vorlage)

WICHTIG: Zusätzliche Boni werden NICHT angezeigt (nur Basis-Mitigationswerte)
"""

import json
import copy
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE


# Maßnahmen-Metadaten (Fokus, Beschreibung und Level-Erklärungen für Nicht-Techniker)
MEASURE_META = {
    "M1": {
        "focus": "Zugriffskontrolle, privilegierte Accounts",
        "description": "Wer darf was im Unternehmen? Diese Maßnahme stellt sicher, dass nur berechtigte Personen Zugang zu sensiblen Systemen haben - wie ein digitaler Türsteher.",
        "levels_simple": {
            1: "Zentrale Benutzerverwaltung, Admins brauchen zwei Faktoren zum Einloggen.",
            2: "Spezielle Verwaltung für Admin-Konten, klare Rollen wer was darf.",
            3: "Zugriff nur bei Bedarf, alle Admin-Aktionen werden aufgezeichnet."
        }
    },
    "M2": {
        "focus": "Sichtbarkeit schaffen, Angriffe erkennen",
        "description": "Wie ein Überwachungssystem für die IT: Alle wichtigen Ereignisse werden protokolliert und ausgewertet, um verdächtige Aktivitäten frühzeitig zu erkennen.",
        "levels_simple": {
            1: "Protokolle werden gesammelt und 30 Tage aufbewahrt.",
            2: "Automatische Alarme bei verdächtigen Mustern (z.B. viele fehlgeschlagene Logins).",
            3: "Externe Sicherheitsexperten überwachen rund um die Uhr."
        }
    },
    "M3": {
        "focus": "Schadsoftware auf Clients/Servern erkennen",
        "description": "Der Virenschutz für alle Computer im Unternehmen. Erkennt und stoppt Schadsoftware, bevor sie Schaden anrichten kann.",
        "levels_simple": {
            1: "Klassischer Virenscanner, der bekannte Schadsoftware erkennt.",
            2: "Erkennt auch unbekannte Bedrohungen anhand von verdächtigem Verhalten, isoliert befallene Rechner automatisch.",
            3: "Vernetzt alle Sicherheitssysteme - erkennt komplexe Angriffe über mehrere Systeme hinweg."
        }
    },
    "M4": {
        "focus": "Daten wiederherstellen, Betrieb aufrechterhalten",
        "description": "Die Lebensversicherung für Unternehmensdaten. Regelmäßige Sicherungskopien ermöglichen es, nach einem Angriff oder Ausfall schnell wieder arbeitsfähig zu sein.",
        "levels_simple": {
            1: "Tägliche Datensicherung, aber ohne regelmäßige Tests ob die Wiederherstellung funktioniert.",
            2: "Backups an mehreren Orten, monatlich wird geprüft ob Wiederherstellung klappt.",
            3: "Backups sind unveränderbar und vom Netzwerk getrennt - selbst Hacker können sie nicht löschen. Schnelle Wiederherstellung garantiert."
        }
    },
    "M5": {
        "focus": "Produktionsnetz vom Office-Netz trennen",
        "description": "Baut digitale Brandmauern zwischen Büro und Produktion. Ein Angriff im Büro soll nicht die Fertigungsanlagen lahmlegen können.",
        "levels_simple": {
            1: "Netzwerke sind logisch getrennt, aber ohne strenge Zugangskontrollen.",
            2: "Echte Firewall zwischen Büro und Produktion - nur explizit erlaubter Datenverkehr kommt durch.",
            3: "Jede Produktionslinie ist einzeln abgesichert, verdächtiger Datenverkehr wird automatisch blockiert."
        }
    },
    "M6": {
        "focus": "Mitarbeitende sensibilisieren, Phishing abwehren",
        "description": "Macht Mitarbeiter zur ersten Verteidigungslinie. Geschulte Mitarbeiter erkennen Betrugsversuche und gefährliche E-Mails, bevor Schaden entsteht.",
        "levels_simple": {
            1: "Einmal im Jahr eine kurze Online-Schulung zum Thema IT-Sicherheit.",
            2: "Regelmäßige Schulungen plus Test-Phishing-Mails, um die Wachsamkeit zu prüfen.",
            3: "Spezielle Trainings je nach Rolle, klarer Prozess zum Melden von Verdachtsfällen."
        }
    },
    "M7": {
        "focus": "Schwachstellen schließen, Angriffsfläche reduzieren",
        "description": "Hält alle Systeme auf dem neuesten Stand. Sicherheitslücken werden durch Updates geschlossen, bevor Angreifer sie ausnutzen können.",
        "levels_simple": {
            1: "Windows-Updates laufen automatisch, Produktionssysteme werden unregelmäßig aktualisiert.",
            2: "Strukturierter monatlicher Update-Prozess, regelmäßige Suche nach Schwachstellen.",
            3: "Kritische Lücken werden priorisiert geschlossen, wöchentliche automatische Schwachstellen-Scans."
        }
    },
    "M8": {
        "focus": "Lieferanten-Risiken managen, Abhängigkeiten absichern",
        "description": "Schützt vor Risiken durch Geschäftspartner. Denn auch über einen gehackten Lieferanten können Angreifer ins Unternehmen gelangen.",
        "levels_simple": {
            1: "Sicherheitsanforderungen stehen im Vertrag, gelegentliche Überprüfungen.",
            2: "Jährliche Sicherheitsbewertung aller wichtigen Lieferanten, Überwachung der Leistungsversprechen.",
            3: "Kontinuierliche Überwachung, Backup-Lieferanten für kritische Teile, Schutz bei Lieferanten-Insolvenz."
        }
    },
    "M9": {
        "focus": "Cloud-Konfigurationen überwachen & absichern",
        "description": "Sichert Daten und Dienste in der Cloud ab. Prüft automatisch, ob Cloud-Systeme sicher konfiguriert sind.",
        "levels_simple": {
            1: "Einfache Überwachung der Cloud-Backups.",
            2: "Automatische Prüfung ob Cloud-Einstellungen den Sicherheitsrichtlinien entsprechen.",
            3: "Umfassende Sicherheit über mehrere Cloud-Anbieter, automatische Prüfung aller Cloud-Konfigurationen."
        }
    },
    "M10": {
        "focus": "Smartphones/Tablets absichern",
        "description": "Schützt mobile Geräte der Mitarbeiter. Verhindert Datenverlust bei Diebstahl und kontrolliert welche Apps genutzt werden dürfen.",
        "levels_simple": {
            1: "PIN-Schutz Pflicht, bei Verlust kann das Gerät aus der Ferne gelöscht werden.",
            2: "Zentrale Verwaltung aller Mobilgeräte, nur zugelassene Apps können installiert werden.",
            3: "Einheitliche Verwaltung aller Endgeräte, aktiver Schutz vor mobilen Bedrohungen."
        }
    }
}

# Event-Metadaten
EVENT_META = {
    "oem_audit": {
        "name": "OEM-Audit",
        "trigger_wave": 2,
        "description": "Der große Automobilkunde prüft die Sicherheitsstandards. Wer gut vorbereitet ist, gewinnt Vertrauen - wer schlecht abschneidet, riskiert Aufträge.",
        "condition": "Ist euer E-Wert >= Zielwert?",
        "effect_positive": "KZ +5",
        "effect_negative": "KZ -3"
    },
    "staff_turnover": {
        "name": "Personalwechsel",
        "trigger_wave": 2,
        "description": "Ohne regelmäßige Schulungen und Sensibilisierung werden IT-Sicherheitsaufgaben zur Belastung. Überlastete Mitarbeiter verlassen das Unternehmen.",
        "condition": "Ist M6 (Security Awareness) < Level 2?",
        "effect_positive": "-",
        "effect_negative": "KZ -2, OPEX +5k€"
    },
    "gdpr_bonus": {
        "name": "DSGVO-Bonus",
        "trigger_wave": 3,
        "description": "Gute Zugriffskontrolle und Logging sind die Basis für Datenschutz-Compliance. Wer hier investiert hat, wird belohnt.",
        "condition": "Sind M1 >= L2 UND M2 >= L2?",
        "effect_positive": "KZ +3, Budget +10k€",
        "effect_negative": "-"
    },
    "investor_confidence": {
        "name": "Investoren-Vertrauen",
        "trigger_wave": 3,
        "description": "Wer viel in Sicherheit investiert, zeigt Weitsicht. Das überzeugt Investoren und Geldgeber.",
        "condition": "Budget-Tier = HIGH?",
        "effect_positive": "KZ +8",
        "effect_negative": "-"
    },
    "compliance_gap": {
        "name": "Compliance-Lücke",
        "trigger_wave": 2,
        "description": "Zu wenig Budget bedeutet Kompromisse bei der Sicherheit. Das fällt spätestens bei Audits und Kundenanfragen negativ auf.",
        "condition": "Budget-Tier = LOW?",
        "effect_positive": "-",
        "effect_negative": "KZ -3"
    }
}


def load_config():
    """Lade Konfiguration aus JSON-Datei"""
    config_path = Path(__file__).parent / "simulation_config.json"
    with open(config_path, "r", encoding="utf-8") as f:
        return json.load(f)


def find_and_replace_text(shape, old_text_part, new_text):
    """Ersetze Text in einem Shape, behält Formatierung bei"""
    if not shape.has_text_frame:
        return False

    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if old_text_part in run.text:
                run.text = run.text.replace(old_text_part, new_text)
                return True
    return False


def set_shape_text(shape, new_text):
    """Setze den kompletten Text eines Shapes"""
    if not shape.has_text_frame:
        return False

    # Behalte die Formatierung des ersten Runs
    if shape.text_frame.paragraphs and shape.text_frame.paragraphs[0].runs:
        first_run = shape.text_frame.paragraphs[0].runs[0]
        first_run.text = new_text
        # Entferne übrige Runs
        for para in shape.text_frame.paragraphs:
            while len(para.runs) > 1:
                para._p.remove(para.runs[-1]._r)
        return True

    # Fallback: Setze Text direkt
    shape.text_frame.paragraphs[0].text = new_text
    return True


def get_dependency_text(measure_data, level):
    """Hole Abhängigkeits-Text für ein bestimmtes Level"""
    dependencies = measure_data.get("dependencies", [])
    for dep in dependencies:
        if dep.get("level") == level:
            reqs = dep.get("requires", [])
            if reqs:
                req = reqs[0]  # Nehme erste Anforderung
                return f"{req['measure']} min. L{req['min_level']}"
    return "-"


def update_massnahmen_template(config, template_path, output_path):
    """
    Aktualisiere das Maßnahmenkarten-Template mit Daten aus der Config.
    Template hat 30 Folien: 10 Maßnahmen × 3 Levels (L1, L2, L3 pro Maßnahme)
    """
    prs = Presentation(template_path)
    measures = config.get("measures", {})

    # Maßnahmen-Reihenfolge
    measure_order = ["M1", "M2", "M3", "M4", "M5", "M6", "M7", "M8", "M9", "M10"]

    slide_idx = 0
    for measure_id in measure_order:
        if measure_id not in measures:
            continue

        measure_data = measures[measure_id]
        meta = MEASURE_META.get(measure_id, {})

        # 3 Levels pro Maßnahme
        for level in [1, 2, 3]:
            if slide_idx >= len(prs.slides):
                print(f"Warnung: Nicht genug Folien im Template für {measure_id} L{level}")
                break

            slide = prs.slides[slide_idx]
            level_data = measure_data["levels"].get(str(level), {})

            # Sammle die Texte zum Ersetzen
            title = f"{measure_id}: {measure_data['name']}"
            subtitle = f"Maßnahme - {meta.get('focus', '')}"
            description = meta.get('description', '')
            level_desc = meta.get('levels_simple', {}).get(level, '')

            cia = level_data.get('cia', {'c': 0, 'i': 0, 'a': 0})
            cia_text = f"C +{cia.get('c', 0)}  |  I +{cia.get('i', 0)}  |  A +{cia.get('a', 0)}"

            init = level_data.get('init', 0)
            opex = level_data.get('opex', 0)
            cost_text = f"Init: {init}k€  |  OPEX: {opex}k€/Welle"

            dep_text = get_dependency_text(measure_data, level)

            level_label = f"Level {level} (L{level})"

            # Durchsuche alle Shapes und aktualisiere Texte
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue

                text = shape.text_frame.text

                # Titel (Mx: Name)
                if text.startswith("M") and ":" in text[:5]:
                    set_shape_text(shape, title)

                # Untertitel (Maßnahme - ...)
                elif text.startswith("Maßnahme"):
                    set_shape_text(shape, subtitle)

                # Beschreibung (längerer Text)
                elif "Wer darf" in text or "Wie ein" in text or "Der Virus" in text or \
                     "Die Lebens" in text or "Baut digit" in text or "Macht Mit" in text or \
                     "Hält alle" in text or "Schützt vor" in text or "Sichert Dat" in text or \
                     "Schützt mob" in text or len(text) > 80:
                    if not any(x in text for x in ["CIA", "Init:", "Level", "Abhäng", "Bedingung"]):
                        set_shape_text(shape, description)

                # Level-Beschreibung (in der Mitte, beschreibender Text)
                elif any(x in text for x in ["Zentrale", "Spezielle", "Zugriff nur", "Protokolle",
                                              "Automatische Alarm", "Externe Sicherheit",
                                              "Klassischer", "Erkennt auch", "Vernetzt alle",
                                              "Tägliche Daten", "Backups an", "Backups sind",
                                              "Netzwerke sind", "Echte Firewall", "Jede Produktion",
                                              "Einmal im", "Regelmäßige Schul", "Spezielle Training",
                                              "Windows-Updates", "Strukturierter", "Kritische Lücken",
                                              "Sicherheitsanford", "Jährliche Sicher", "Kontinuierliche",
                                              "Einfache Über", "Automatische Prüf", "Umfassende Sicher",
                                              "PIN-Schutz", "Zentrale Verwalt", "Einheitliche Verwalt"]):
                    set_shape_text(shape, level_desc)

                # Level-Label
                elif text.startswith("Level"):
                    set_shape_text(shape, level_label)

                # CIA-Beitrag Wert
                elif "C +" in text and "|" in text:
                    set_shape_text(shape, cia_text)

                # Kosten Wert
                elif "Init:" in text and "OPEX:" in text:
                    set_shape_text(shape, cost_text)

                # Abhängigkeiten Wert
                elif text.strip() == "-" or (text.startswith("M") and "min." in text):
                    if "Abhängig" not in text:
                        set_shape_text(shape, dep_text)

            slide_idx += 1

    # Speichern
    prs.save(output_path)
    print(f"Maßnahmenkarten aus Template generiert: {output_path}")


def duplicate_slide(prs, slide_index):
    """Dupliziere eine Folie"""
    template_slide = prs.slides[slide_index]

    # Füge neue Folie mit gleichem Layout hinzu
    slide_layout = template_slide.slide_layout
    new_slide = prs.slides.add_slide(slide_layout)

    # Kopiere alle Shapes
    for shape in template_slide.shapes:
        # Überspringe Platzhalter, die automatisch kopiert werden
        if shape.is_placeholder:
            # Finde entsprechenden Platzhalter in neuer Folie
            for new_shape in new_slide.shapes:
                if new_shape.is_placeholder and new_shape.placeholder_format.idx == shape.placeholder_format.idx:
                    if shape.has_text_frame:
                        # Kopiere Text
                        new_shape.text_frame.clear()
                        for para_idx, para in enumerate(shape.text_frame.paragraphs):
                            if para_idx == 0:
                                new_para = new_shape.text_frame.paragraphs[0]
                            else:
                                new_para = new_shape.text_frame.add_paragraph()
                            new_para.text = para.text
                    break

    return new_slide


def copy_slide_shapes(source_slide, target_slide):
    """Kopiere alle Nicht-Platzhalter-Shapes von einer Folie zur anderen"""
    from pptx.util import Emu
    from copy import deepcopy
    import lxml.etree as etree

    for shape in source_slide.shapes:
        # Überspringe Platzhalter (werden automatisch vom Layout erstellt)
        if shape.is_placeholder:
            continue

        # Kopiere das XML-Element
        el = shape.element
        new_el = deepcopy(el)
        target_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')


def update_events_template(config, template_path, output_path):
    """
    Aktualisiere das Events-Template mit Daten aus der Config.
    Template hat 1 Folie als Vorlage, die für jedes Event dupliziert wird.

    Platzhalter werden nach Index identifiziert:
    - idx 13: Untertitel (Event – Nach Welle X)
    - idx 0: Titel (Event-Name)
    - idx 14: Beschreibung (1. Spalte oben) - NUR diese verwenden
    - idx 17: Bedingung (1. Spalte unten) - NUR diese verwenden
    - idx 18-21: Weitere Spalten - leer lassen

    Nicht-Platzhalter (TEXT_BOX):
    - top > 4.5", left < 4": Positiver Effekt
    - top > 4.5", left > 4": Negativer Effekt
    """
    from copy import deepcopy

    prs = Presentation(template_path)

    # Event-Reihenfolge
    event_order = ["oem_audit", "staff_turnover", "gdpr_bonus", "investor_confidence", "compliance_gap"]

    # Erste Folie als Vorlage
    template_slide = prs.slides[0]

    for event_idx, event_id in enumerate(event_order):
        meta = EVENT_META.get(event_id, {})

        if event_idx == 0:
            # Erste Folie: Aktualisiere direkt
            slide = template_slide
        else:
            # Weitere Folien: Dupliziere Vorlage via XML-Kopie
            slide_layout = template_slide.slide_layout
            new_slide = prs.slides.add_slide(slide_layout)

            # Kopiere alle Nicht-Platzhalter-Shapes via XML
            for shape in template_slide.shapes:
                if not shape.is_placeholder:
                    el = shape.element
                    new_el = deepcopy(el)
                    new_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')

            slide = new_slide

        # Texte zum Ersetzen
        title = meta.get('name', event_id)
        subtitle = f"Event – Nach Welle {meta.get('trigger_wave', '?')}"
        description = meta.get('description', '')
        condition = meta.get('condition', '')
        effect_pos = meta.get('effect_positive', '-')
        effect_neg = meta.get('effect_negative', '-')

        # Aktualisiere Texte in der Folie
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            top = shape.top.inches
            left = shape.left.inches
            text = shape.text_frame.text

            # Platzhalter basierend auf Index identifizieren
            if shape.is_placeholder:
                idx = shape.placeholder_format.idx

                # Untertitel (idx 13)
                if idx == 13:
                    set_shape_text(shape, subtitle)

                # Titel (idx 0)
                elif idx == 0:
                    set_shape_text(shape, title)

                # Beschreibung (idx 14 - erste Spalte oben)
                elif idx == 14:
                    set_shape_text(shape, description)

                # Bedingung (idx 17 - erste Spalte unten)
                elif idx == 17:
                    set_shape_text(shape, condition)

                # Weitere Spalten leer lassen (idx 18-21)
                elif idx in [18, 19, 20, 21]:
                    set_shape_text(shape, "")

            # Nicht-Platzhalter (TEXT_BOX für Effekte)
            else:
                # Labels nicht ändern
                if "erfüllt" in text:
                    continue

                # Effekte (unten, top > 4.5")
                if top >= 4.5:
                    # Positiver Effekt (links, left < 4")
                    if left < 4:
                        set_shape_text(shape, effect_pos if effect_pos else "-")

                    # Negativer Effekt (rechts, left > 4")
                    elif left > 4:
                        set_shape_text(shape, effect_neg if effect_neg else "-")

    # Speichern
    prs.save(output_path)
    print(f"Event-Karten aus Template generiert: {output_path}")


def main():
    """Hauptfunktion"""
    # Konfiguration laden
    config = load_config()

    # Pfade
    base_dir = Path(__file__).parent
    template_dir = base_dir / "pptx_output"
    output_dir = base_dir / "pptx_output"

    # Template-Pfade
    massnahmen_template = template_dir / "Massnahmenkarten_Security-Game.pptx"
    events_template = template_dir / "Events_Security-Game.pptx"

    # Output-Pfade
    massnahmen_output = output_dir / "Massnahmenkarten_Generated.pptx"
    events_output = output_dir / "Events_Generated.pptx"

    # Prüfe ob Templates existieren
    if not massnahmen_template.exists():
        print(f"FEHLER: Template nicht gefunden: {massnahmen_template}")
        return

    if not events_template.exists():
        print(f"FEHLER: Template nicht gefunden: {events_template}")
        return

    # Generiere aus Templates
    print("Generiere PowerPoint-Dateien aus Templates...")
    print()

    update_massnahmen_template(config, massnahmen_template, massnahmen_output)
    update_events_template(config, events_template, events_output)

    print()
    print("Fertig! Generierte Dateien:")
    print(f"  - {massnahmen_output}")
    print(f"  - {events_output}")


if __name__ == "__main__":
    main()
