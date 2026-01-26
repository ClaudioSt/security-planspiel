#!/usr/bin/env python3
"""
PowerPoint-Generator für Maßnahmenkarten und Events
Security-Planspiel MechTech

Erstellt:
- Maßnahmenkarten (M1-M10) als PowerPoint-Folien
- Event-Karten als PowerPoint-Folien

WICHTIG: Zusätzliche Boni werden NICHT angezeigt (nur Basis-Mitigationswerte)
"""

import json
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


# Farbschema
COLORS = {
    "primary": RGBColor(0, 82, 147),      # Dunkelblau
    "secondary": RGBColor(0, 122, 194),   # Mittelblau
    "accent": RGBColor(255, 153, 0),      # Orange
    "success": RGBColor(0, 153, 76),      # Grün
    "warning": RGBColor(255, 193, 7),     # Gelb
    "danger": RGBColor(220, 53, 69),      # Rot
    "light": RGBColor(248, 249, 250),     # Hellgrau
    "dark": RGBColor(52, 58, 64),         # Dunkelgrau
    "white": RGBColor(255, 255, 255),
    "c_color": RGBColor(0, 123, 255),     # Blau für Confidentiality
    "i_color": RGBColor(255, 193, 7),     # Gelb für Integrity
    "a_color": RGBColor(40, 167, 69),     # Grün für Availability
}

# Maßnahmen-Metadaten (Fokus und Beschreibung)
MEASURE_META = {
    "M1": {
        "focus": "Zugriffskontrolle, privilegierte Accounts",
        "levels_desc": {
            1: "Zentrale AD, MFA für Admins, Passwortrichtlinie",
            2: "PAM (Privileged Access Management), Rollenkonzept",
            3: "JIT-Access, Session-Recording, Access-Reviews"
        }
    },
    "M2": {
        "focus": "Sichtbarkeit schaffen, Angriffe erkennen",
        "levels_desc": {
            1: "Windows Event Logs zentral sammeln, 30 Tage",
            2: "SIEM mit Use Cases (Failed Logins, Priv. Escalation)",
            3: "MDR (Managed Detection & Response) mit 24/7-SOC"
        }
    },
    "M3": {
        "focus": "Schadsoftware auf Clients/Servern erkennen",
        "levels_desc": {
            1: "Aktueller Virenscanner, automatische Updates",
            2: "EDR mit Behavioral Analysis, Auto-Isolation",
            3: "XDR (Korrelation Endpoints, Netzwerk, Cloud)"
        }
    },
    "M4": {
        "focus": "Daten wiederherstellen, Betrieb aufrechterhalten",
        "levels_desc": {
            1: "Tägliche Backups Office-IT, 14 Tage, kein Test",
            2: "3-2-1-Regel, monatliche Restore-Tests",
            3: "Immutable Backups (Air-Gap), DR-Plan, RTO <4h"
        }
    },
    "M5": {
        "focus": "Produktionsnetz vom Office-Netz trennen",
        "levels_desc": {
            1: "Logische VLANs, keine Firewall-Regeln",
            2: "Firewall zwischen OT/IT, Whitelist-Prinzip",
            3: "IDS/IPS in OT, Micro-Segmentierung (Zonen/Linie)"
        }
    },
    "M6": {
        "focus": "Mitarbeitende sensibilisieren, Phishing abwehren",
        "levels_desc": {
            1: "Jährliche Pflicht-Schulung (E-Learning, 30 Min)",
            2: "Quartalsweise Schulungen + Phishing-Simulationen",
            3: "Rollenspez. Training, Incident-Meldeprozess etabliert"
        }
    },
    "M7": {
        "focus": "Schwachstellen schließen, Angriffsfläche reduzieren",
        "levels_desc": {
            1: "Automatische Windows-Updates, sporadisch OT-Patches",
            2: "Monatliche Patch-Zyklen, Vulnerability-Scans (quartal)",
            3: "Risk-based Patching (CVSS), wöchentliche Scans, Tool"
        }
    },
    "M8": {
        "focus": "Lieferanten-Risiken managen, Abhängigkeiten absichern",
        "levels_desc": {
            1: "Vertragliche Security-Klauseln, Ad-hoc-Checks",
            2: "Supplier-Assessments (jährlich), SLA-Monitoring",
            3: "Continuous Monitoring, Dual-Sourcing, Escrow-Verträge"
        }
    },
    "M9": {
        "focus": "Cloud-Konfigurationen überwachen & absichern",
        "levels_desc": {
            1: "Basis-Monitoring Cloud-Backup (Veeam)",
            2: "CSPM-Tool (z.B. Prisma Cloud), Compliance-Checks",
            3: "Multi-Cloud-Security, IaC-Scanning, CNAPP"
        }
    },
    "M10": {
        "focus": "Smartphones/Tablets absichern",
        "levels_desc": {
            1: "PIN-Pflicht, Remote-Wipe bei Verlust",
            2: "MDM-Plattform (Intune/Jamf), App-Whitelisting",
            3: "UEM (Unified Endpoint Mgmt.), MTD (Mobile Threat Def.)"
        }
    }
}

# Event-Metadaten mit deutschen Beschreibungen
EVENT_META = {
    "oem_audit": {
        "name": "OEM-Audit",
        "icon": "📋",
        "condition": "Nach Welle 1: E-Wert ≥ Zielwert?",
        "effect_positive": "KZ +5",
        "effect_negative": "KZ -3"
    },
    "staff_turnover": {
        "name": "Personalwechsel",
        "icon": "👥",
        "condition": "M6 (Security Awareness) < Level 2?",
        "effect_positive": None,
        "effect_negative": "KZ -2, OPEX +5k€"
    },
    "gdpr_bonus": {
        "name": "DSGVO-Bonus",
        "icon": "🏆",
        "condition": "M1 ≥ L2 UND M2 ≥ L2?",
        "effect_positive": "KZ +3, Budget +10k€",
        "effect_negative": None
    },
    "investor_confidence": {
        "name": "Investoren-Vertrauen",
        "icon": "💰",
        "condition": "Budget-Tier = HIGH",
        "effect_positive": "KZ +8",
        "effect_negative": None
    },
    "compliance_gap": {
        "name": "Compliance-Lücke",
        "icon": "⚠️",
        "condition": "Budget-Tier = LOW",
        "effect_positive": None,
        "effect_negative": "KZ -3"
    }
}

# Angriffsnamen für Mitigation-Anzeige
ATTACK_NAMES = {
    "ransomware": "Ransomware",
    "ot": "OT-Störung",
    "exfil": "Exfiltration"
}


def load_config():
    """Lade Konfiguration aus JSON-Datei"""
    config_path = Path(__file__).parent / "simulation_config.json"
    with open(config_path, "r", encoding="utf-8") as f:
        return json.load(f)


def add_title_slide(prs, title, subtitle=""):
    """Füge Titelfolie hinzu"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Hintergrund
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = COLORS["primary"]
    background.line.fill.background()

    # Titel
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.5), Inches(9), Inches(1.5)
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLORS["white"]
    p.alignment = PP_ALIGN.CENTER

    # Untertitel
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(4), Inches(9), Inches(1)
        )
        tf = subtitle_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = COLORS["light"]
        p.alignment = PP_ALIGN.CENTER

    return slide


def add_measure_slide(prs, measure_id, measure_data, meta):
    """Füge Maßnahmenkarten-Folie hinzu"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Header-Bereich (blau)
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.4)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = COLORS["primary"]
    header.line.fill.background()

    # Titel
    title_box = slide.shapes.add_textbox(
        Inches(0.3), Inches(0.2), Inches(9.4), Inches(0.6)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"{measure_id}: {measure_data['name']}"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLORS["white"]

    # Fokus-Zeile
    focus_box = slide.shapes.add_textbox(
        Inches(0.3), Inches(0.85), Inches(9.4), Inches(0.4)
    )
    tf = focus_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"Fokus: {meta['focus']}"
    p.font.size = Pt(16)
    p.font.italic = True
    p.font.color.rgb = COLORS["light"]

    # Drei Level-Spalten
    levels = [1, 2, 3]
    col_width = Inches(3.1)
    col_start = Inches(0.25)
    col_gap = Inches(0.1)

    level_colors = [COLORS["success"], COLORS["warning"], COLORS["danger"]]
    level_names = ["BASIS", "STANDARD", "ERWEITERT"]

    for i, level in enumerate(levels):
        level_data = measure_data["levels"].get(str(level), {})
        x_pos = col_start + i * (col_width + col_gap)

        # Level-Box
        level_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x_pos, Inches(1.6), col_width, Inches(5.5)
        )
        level_box.fill.solid()
        level_box.fill.fore_color.rgb = COLORS["white"]
        level_box.line.color.rgb = level_colors[i]
        level_box.line.width = Pt(3)

        # Level-Header
        header_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, x_pos, Inches(1.6), col_width, Inches(0.5)
        )
        header_shape.fill.solid()
        header_shape.fill.fore_color.rgb = level_colors[i]
        header_shape.line.fill.background()

        level_title = slide.shapes.add_textbox(x_pos, Inches(1.65), col_width, Inches(0.4))
        tf = level_title.text_frame
        p = tf.paragraphs[0]
        p.text = f"LEVEL {level} ({level_names[i]})"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = COLORS["white"]
        p.alignment = PP_ALIGN.CENTER

        # Beschreibung
        desc_box = slide.shapes.add_textbox(
            x_pos + Inches(0.1), Inches(2.15), col_width - Inches(0.2), Inches(0.8)
        )
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = meta["levels_desc"].get(level, "")
        p.font.size = Pt(11)
        p.font.color.rgb = COLORS["dark"]

        # CIA-Beitrag
        cia = level_data.get("cia", {"c": 0, "i": 0, "a": 0})
        cia_box = slide.shapes.add_textbox(
            x_pos + Inches(0.1), Inches(2.95), col_width - Inches(0.2), Inches(0.6)
        )
        tf = cia_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "CIA-Beitrag:"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = COLORS["dark"]

        p = tf.add_paragraph()
        p.text = f"C +{cia.get('c', 0)}  |  I +{cia.get('i', 0)}  |  A +{cia.get('a', 0)}"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = COLORS["secondary"]

        # Kosten
        init = level_data.get("init", 0)
        opex = level_data.get("opex", 0)
        cost_box = slide.shapes.add_textbox(
            x_pos + Inches(0.1), Inches(3.55), col_width - Inches(0.2), Inches(0.6)
        )
        tf = cost_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "Kosten:"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = COLORS["dark"]

        p = tf.add_paragraph()
        p.text = f"Init: {init}k€  |  OPEX: {opex}k€/Welle"
        p.font.size = Pt(11)
        p.font.color.rgb = COLORS["dark"]

        # Mitigation (OHNE Bonus - nur Basiswerte!)
        mitigation = level_data.get("mitigation", {})
        mit_box = slide.shapes.add_textbox(
            x_pos + Inches(0.1), Inches(4.15), col_width - Inches(0.2), Inches(1.4)
        )
        tf = mit_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "Mitigation:"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = COLORS["dark"]

        has_mitigation = False
        for attack_key, attack_name in ATTACK_NAMES.items():
            mit_value = mitigation.get(attack_key, 0)
            if mit_value != 0:
                has_mitigation = True
                p = tf.add_paragraph()
                p.text = f"• {attack_name}: {mit_value}"
                p.font.size = Pt(11)
                p.font.color.rgb = COLORS["success"] if mit_value < 0 else COLORS["dark"]

        if not has_mitigation:
            p = tf.add_paragraph()
            p.text = "• Keine direkte Wirkung"
            p.font.size = Pt(11)
            p.font.color.rgb = COLORS["dark"]
            p.font.italic = True

        # Recovery (nur für M4)
        recovery = level_data.get("recovery", 0)
        if recovery > 0:
            p = tf.add_paragraph()
            p.text = f"• Recovery: {int(recovery * 100)}%"
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = COLORS["accent"]

        # Abhängigkeiten
        if measure_data.get("dependencies"):
            dep_box = slide.shapes.add_textbox(
                x_pos + Inches(0.1), Inches(5.5), col_width - Inches(0.2), Inches(0.5)
            )
            tf = dep_box.text_frame
            tf.word_wrap = True

            for dep in measure_data["dependencies"]:
                if dep.get("level") == level:
                    for req in dep.get("requires", []):
                        p = tf.paragraphs[0]
                        p.text = f"Benötigt: {req['measure']} min. L{req['min_level']}"
                        p.font.size = Pt(10)
                        p.font.italic = True
                        p.font.color.rgb = COLORS["danger"]

    return slide


def add_event_slide(prs, event_id, event_data, meta):
    """Füge Event-Folie hinzu"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Header-Bereich (orange)
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.8)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = COLORS["accent"]
    header.line.fill.background()

    # Event-Titel
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(9), Inches(0.8)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"{meta['icon']} {meta['name']}"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLORS["white"]
    p.alignment = PP_ALIGN.CENTER

    # Trigger-Welle
    trigger_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.15), Inches(9), Inches(0.5)
    )
    tf = trigger_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"Trigger: Nach Welle {event_data.get('trigger_wave', '?')}"
    p.font.size = Pt(20)
    p.font.color.rgb = COLORS["white"]
    p.alignment = PP_ALIGN.CENTER

    # Beschreibung
    desc_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.2), Inches(9), Inches(1)
    )
    tf = desc_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = event_data.get("description", "")
    p.font.size = Pt(22)
    p.font.color.rgb = COLORS["dark"]
    p.alignment = PP_ALIGN.CENTER

    # Bedingung
    cond_shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(3.4), Inches(9), Inches(1.2)
    )
    cond_shape.fill.solid()
    cond_shape.fill.fore_color.rgb = COLORS["light"]
    cond_shape.line.color.rgb = COLORS["secondary"]
    cond_shape.line.width = Pt(2)

    cond_box = slide.shapes.add_textbox(
        Inches(0.7), Inches(3.5), Inches(8.6), Inches(1)
    )
    tf = cond_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Bedingung:"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLORS["dark"]

    p = tf.add_paragraph()
    p.text = meta["condition"]
    p.font.size = Pt(18)
    p.font.color.rgb = COLORS["primary"]

    # Effekte
    y_pos = Inches(4.9)

    # Positiver Effekt
    if meta["effect_positive"]:
        pos_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), y_pos, Inches(4.3), Inches(1.2)
        )
        pos_shape.fill.solid()
        pos_shape.fill.fore_color.rgb = COLORS["success"]
        pos_shape.line.fill.background()

        pos_box = slide.shapes.add_textbox(
            Inches(0.7), y_pos + Inches(0.15), Inches(3.9), Inches(0.9)
        )
        tf = pos_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "Bei Erfüllung:"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = COLORS["white"]

        p = tf.add_paragraph()
        p.text = meta["effect_positive"]
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = COLORS["white"]

    # Negativer Effekt
    if meta["effect_negative"]:
        neg_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.2), y_pos, Inches(4.3), Inches(1.2)
        )
        neg_shape.fill.solid()
        neg_shape.fill.fore_color.rgb = COLORS["danger"]
        neg_shape.line.fill.background()

        neg_box = slide.shapes.add_textbox(
            Inches(5.4), y_pos + Inches(0.15), Inches(3.9), Inches(0.9)
        )
        tf = neg_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "Bei Nicht-Erfüllung:"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = COLORS["white"]

        p = tf.add_paragraph()
        p.text = meta["effect_negative"]
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = COLORS["white"]

    # Wenn nur ein Effekt vorhanden
    if meta["effect_positive"] and not meta["effect_negative"]:
        # Nur positiver Effekt - zentrieren
        pass  # Layout ist schon okay
    elif meta["effect_negative"] and not meta["effect_positive"]:
        # Nur negativer Effekt - zentrieren
        pass  # Layout ist schon okay

    return slide


def generate_massnahmen_pptx(config, output_path):
    """Generiere PowerPoint mit allen Maßnahmenkarten"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Titelfolie
    add_title_slide(
        prs,
        "Maßnahmenkarten",
        "Security-Planspiel MechTech\n10 Maßnahmen × 3 Level"
    )

    # Maßnahmen-Folien
    measures = config.get("measures", {})
    for measure_id in sorted(measures.keys(), key=lambda x: int(x[1:])):
        measure_data = measures[measure_id]
        meta = MEASURE_META.get(measure_id, {
            "focus": measure_data.get("name", ""),
            "levels_desc": {1: "", 2: "", 3: ""}
        })
        add_measure_slide(prs, measure_id, measure_data, meta)

    # Speichern
    prs.save(output_path)
    print(f"Maßnahmenkarten gespeichert: {output_path}")


def generate_events_pptx(config, output_path):
    """Generiere PowerPoint mit allen Events"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Titelfolie
    add_title_slide(
        prs,
        "Event-Karten",
        "Security-Planspiel MechTech\nDynamische Spielereignisse"
    )

    # Event-Folien
    events = config.get("events", {})
    event_order = ["oem_audit", "staff_turnover", "gdpr_bonus", "investor_confidence", "compliance_gap"]

    for event_id in event_order:
        if event_id in events:
            event_data = events[event_id]
            meta = EVENT_META.get(event_id, {
                "name": event_id,
                "icon": "📌",
                "condition": "",
                "effect_positive": None,
                "effect_negative": None
            })
            add_event_slide(prs, event_id, event_data, meta)

    # Speichern
    prs.save(output_path)
    print(f"Event-Karten gespeichert: {output_path}")


def main():
    """Hauptfunktion"""
    # Konfiguration laden
    config = load_config()

    # Output-Verzeichnis
    output_dir = Path(__file__).parent / "pptx_output"
    output_dir.mkdir(exist_ok=True)

    # PowerPoint-Dateien generieren
    generate_massnahmen_pptx(config, output_dir / "Massnahmenkarten.pptx")
    generate_events_pptx(config, output_dir / "Events.pptx")

    print("\nPowerPoint-Dateien erfolgreich erstellt!")
    print(f"Ausgabeverzeichnis: {output_dir}")


if __name__ == "__main__":
    main()
